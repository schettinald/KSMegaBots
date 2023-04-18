import io
from azure.storage.blob import BlobServiceClient, BlobClient, ContainerClient
import chardet
from docx import Document
from pptx import Presentation
import PyPDF2

# Functions to read different file types
def read_word_document(content_bytes):
    with io.BytesIO(content_bytes) as byte_stream:
        doc = Document(byte_stream)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

def read_powerpoint_document(content_bytes):
    with io.BytesIO(content_bytes) as byte_stream:
        ppt = Presentation(byte_stream)
        text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if shape.has_text_frame])
    return text

def read_pdf_document(content_bytes):
    with io.BytesIO(content_bytes) as byte_stream:
        reader = PyPDF2.PdfFileReader(byte_stream)
        text = "\n".join([reader.getPage(i).extractText() for i in range(reader.getNumPages())])
    return text

# Replace these with your storage account name, container name, and access key
storage_account_name = "Storage Account Name" 
container_name = "Container Name" 
access_key = "Access Key" 

 

# Connect to the Blob Storage account
connection_string = f"DefaultEndpointsProtocol=https;AccountName={storage_account_name};AccountKey={access_key};EndpointSuffix=core.windows.net"
blob_service_client = BlobServiceClient.from_connection_string(connection_string)

# Connect to the container
container_client = blob_service_client.get_container_client(container_name)

# Get a list of blobs (files) in the container
blobs = container_client.list_blobs()

# Read the content of each blob and write it to the output file
with open("output.txt", "w", encoding="utf-8") as output_file:
    for blob in blobs:
        blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob.name)
        content_bytes = blob_client.download_blob().readall()

        if blob.name.lower().endswith('.docx'):
            content = read_word_document(content_bytes)
        elif blob.name.lower().endswith('.pptx'):
            content = read_powerpoint_document(content_bytes)
        elif blob.name.lower().endswith('.pdf'):
            content = read_pdf_document(content_bytes)
        else:
            detected_encoding = chardet.detect(content_bytes)["encoding"]
            if detected_encoding is None:
                detected_encoding = "utf-8"
            content = content_bytes.decode(detected_encoding)
        
        # Print the content
        print(f"Blob {blob.name} content:\n{content}\n")

        # Write the content to the output file
        print(f"Blob {blob.name} content:\n{content}\n", file=output_file)
