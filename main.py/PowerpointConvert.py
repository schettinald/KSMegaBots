import sys
import collections 
import collections.abc
from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    text.append(paragraph.text)

    return "\n".join(text)

if __name__ == "__main__":
    pptx_file = sys.argv[1]
    output_file = sys.argv[2]

    text = extract_text_from_pptx(pptx_file)

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(text)

    print(f"Text extracted from {pptx_file} and saved to {output_file}")
